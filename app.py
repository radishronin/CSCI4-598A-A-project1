""" A Flask application for LLM interaction. """
import io
import os
import re
from flask import Flask, request, jsonify, render_template
from langgraph.graph.state import CompiledStateGraph
from langgraph.prebuilt import create_react_agent
from langchain_core.language_models import BaseChatModel
from langchain_core.messages import AIMessage
from langchain_google_genai import ChatGoogleGenerativeAI
from langchain_openai import ChatOpenAI
from pypdf import PdfReader
from docx import Document
from pptx import Presentation
from bs4 import BeautifulSoup

app = Flask(__name__)

''' Constants '''
GEMINI_MODEL = "gemini-2.0-flash"
OPENAI_MODEL = "gpt-4o-mini"
MAX_CONTEXT_CHARS = 20_000

LAST_DOC_TEXT = ""


def extract_text(name: str, data: bytes) -> str:
    """Extract plain text from supported document formats."""
    ext = os.path.splitext(name or "")[1].lower()
    text_chunks: list[str] = []

    try:
        if ext == ".pdf":
            reader = PdfReader(io.BytesIO(data))
            for page in reader.pages:
                try:
                    page_text = page.extract_text() or ""
                except Exception:
                    page_text = ""
                if page_text:
                    text_chunks.append(page_text)
        elif ext == ".docx":
            document = Document(io.BytesIO(data))
            text_chunks.extend(paragraph.text for paragraph in document.paragraphs if paragraph.text)
        elif ext == ".pptx":
            presentation = Presentation(io.BytesIO(data))
            for slide in presentation.slides:
                for shape in slide.shapes:
                    text = getattr(shape, "text", "")
                    if text:
                        text_chunks.append(text)
        elif ext in {".html", ".htm"}:
            soup = BeautifulSoup(data, "lxml")
            text = soup.get_text(separator=" ")
            if text:
                text_chunks.append(text)
        elif ext in {".txt", ".md"}:
            text = data.decode("utf-8", "ignore")
            if text:
                text_chunks.append(text)
        else:
            text = data.decode("utf-8", "ignore")
            if text:
                text_chunks.append(text)
    except Exception:
        return ""

    normalized = re.sub(r"\s+", " ", " ".join(text_chunks)).strip()
    return normalized

def get_environment_api_key(llm_choice: str) -> str:
    """Get the API key for the selected LLM from ./api_keys/{llm_choice}"""
    key_dir = os.path.join(".", "api_keys")
    key_path = os.path.join(key_dir, llm_choice)
    if not os.path.exists(key_path):
        return ""
    try:
        with open(key_path, "r", encoding="utf-8") as f:
            value = f.read().strip()
            return value if value else ""
    except OSError:
        return ""

def set_environment_api_key(llm_choice: str, api_key: str):
    """Persist the API key for the selected LLM into ./api_keys/{llm_choice}"""
    key_dir = os.path.join(".", "api_keys")
    os.makedirs(key_dir, exist_ok=True)
    key_path = os.path.join(key_dir, llm_choice)
    # Overwrite any existing value with the new key
    with open(key_path, "w", encoding="utf-8") as f:
        f.write(api_key.strip())

@app.route("/")
def index():
    """ Render the index.html template """
    return render_template("index.html")

def _message_to_text(msg) -> str:
    """Extract plain text from message content which can be a str or list of parts."""
    try:
        content = getattr(msg, "content", msg)
        if isinstance(content, str):
            return content
        if isinstance(content, list):
            # Concatenate text-like parts
            parts = []
            for p in content:
                # parts may be dicts or objects with 'type'/'text' fields
                if isinstance(p, dict):
                    t = p.get("text") or p.get("content") or ""
                    if isinstance(t, str):
                        parts.append(t)
                else:
                    t = getattr(p, "text", None)
                    if isinstance(t, str):
                        parts.append(t)
            return "".join(parts)
        return str(content)
    except Exception:
        return ""

@app.route("/api/prompt", methods=["POST"])
def receive_prompt():
    """ 
    1. Retrieve frontend variables
    2. Create Langchain agent, return response 
    """
    data = request.get_json(silent=True) or {}

    prompt_text = (data.get("prompt") or "").strip()
    if not prompt_text:
        return jsonify({"error": "missing fields: prompt"}), 400

    provider = (data.get("provider") or data.get("llm_choice") or "openai").strip().lower() or "openai"

    if provider not in {"openai", "gemini"}:
        return jsonify({"error": f"unsupported provider: {provider}"}), 400

    request_api_key = data.get("apiKey") or data.get("api_key") or ""
    stored_api_key = get_environment_api_key(provider)

    global LAST_DOC_TEXT
    prompt_with_context = prompt_text
    if LAST_DOC_TEXT:
        prompt_with_context = (
            "Use the following context to answer. If irrelevant, say so briefly.\n"
            f"{LAST_DOC_TEXT}\n\nQuestion: {prompt_text}"
        )

    if provider == "openai":
        api_key = os.getenv("OPENAI_API_KEY") or stored_api_key or request_api_key
    else:
        api_key = os.getenv("GEMINI_API_KEY") or os.getenv("GOOGLE_API_KEY") or stored_api_key or request_api_key

    if not api_key:
        return jsonify({"error": f"missing API key for provider: {provider}"}), 400

    try:
        if provider == "gemini":
            langchain_llm: BaseChatModel = ChatGoogleGenerativeAI(
                model=GEMINI_MODEL,
                google_api_key=api_key
            )
        else:
            langchain_llm = ChatOpenAI(
                model=OPENAI_MODEL,
                api_key=api_key
            )

        langchain_agent: CompiledStateGraph = create_react_agent(
            model=langchain_llm,
            tools=[]
        )

        output_chunks: list[str] = []
        for step in langchain_agent.stream({"messages": [prompt_with_context]}, stream_mode="values"):
            msg = step["messages"][-1]
            if isinstance(msg, AIMessage):
                text = _message_to_text(msg)
                if text:
                    output_chunks.append(text)

        output_text = "".join(output_chunks).strip()
    except Exception as exc:  # pragma: no cover - best effort error handling
        return jsonify({"error": str(exc), "provider": provider}), 500

    return jsonify({"output": output_text, "provider": provider})

@app.route("/api/set-api-key", methods=["POST"])
def set_api_key():
    """ Set the API key for the selected LLM """
    data = request.get_json(silent=True) or {}
    llm_choice = data.get("llm_choice", "")
    api_key = data.get("api_key", "")
  
    if llm_choice == "":
        return jsonify({"ok": False, "error": "SNO: no LLM choice provided."}), 400
        
    if api_key == "":
        return jsonify({"ok": False, "error": "SNO: no API key provided."}), 400
    
    # Set the environment variable based on LLM choice
    set_environment_api_key(llm_choice, api_key)
    
    return jsonify({"ok": True, "message": f"API key set for {llm_choice}"})


@app.route("/api/upload-files", methods=["POST"])
def upload_files():
    """Handle multipart file uploads and load text into the context buffer."""
    file_objects = list(request.files.getlist("files"))
    alt_files = request.files.getlist("files[]")
    if alt_files:
        file_objects.extend(alt_files)

    if not file_objects:
        return jsonify({"error": "no files provided"}), 400

    global LAST_DOC_TEXT
    new_chunks: list[str] = []

    for storage in file_objects:
        try:
            file_bytes = storage.read() or b""
        except Exception:
            file_bytes = b""
        text = extract_text(storage.filename or "", file_bytes)
        if text:
            new_chunks.append(text)

    if new_chunks:
        buffer_text = "\n\n".join(new_chunks)
        if LAST_DOC_TEXT:
            buffer_text = f"{LAST_DOC_TEXT}\n\n{buffer_text}"
        LAST_DOC_TEXT = buffer_text[-MAX_CONTEXT_CHARS:]

    return jsonify({"ok": True, "chars": len(LAST_DOC_TEXT)})


@app.route("/api/context-chars", methods=["GET"])
def get_context_chars():
    """Return the number of characters currently stored in the context buffer."""
    return jsonify({"chars": len(LAST_DOC_TEXT)})


@app.route("/api/clear-context", methods=["POST"])
def clear_context():
    """Clear the in-memory context buffer."""
    global LAST_DOC_TEXT
    LAST_DOC_TEXT = ""
    return jsonify({"ok": True})

if __name__ == "__main__":
    app.run(debug=True)
