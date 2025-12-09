"""Blueprint handling RAG-related endpoints and utilities."""
import base64
import os
from io import BytesIO
from pathlib import Path
import logging
import sys
import time
import json
import traceback

import pandas as pd
import pdfplumber
from docx import Document as DocxDocument
from pptx import Presentation
from bs4 import BeautifulSoup
import urllib.request
import urllib.error
from flask import Blueprint, Response, jsonify, request, stream_with_context, render_template, redirect, url_for
from langchain_core.language_models import BaseChatModel
from langchain_core.messages import AIMessage
from langchain_google_genai import ChatGoogleGenerativeAI
from langgraph.graph.state import CompiledStateGraph
from langgraph.prebuilt import create_react_agent
from llama_index.core import Document, Settings, StorageContext
from llama_index.core import VectorStoreIndex, load_index_from_storage
from llama_index.core.langchain_helpers.agents import IndexToolConfig, LlamaIndexTool
from llama_index.core.node_parser import SentenceSplitter
from llama_index.core.vector_stores.simple import SimpleVectorStore
from llama_index.embeddings.google_genai import GoogleGenAIEmbedding
from llama_index.llms.google_genai import GoogleGenAI

rag_bp = Blueprint("rag", __name__, url_prefix = "/rag")

# Constants
GEMINI_MODEL = "gemini-2.0-flash"
GEMINI_EMBEDDING_MODEL = "text-embedding-004"
INDEX_PATH = "./rag_documents"
API_KEY_PATH = "./api_keys"
is_first_llm_run = 1
is_first_embed_run = 1
# Notes storage file (simple JSON store)
BASE_DIR = Path(__file__).resolve().parent.parent
NOTES_DIR = BASE_DIR / "resources"
NOTES_FILE = NOTES_DIR / "notes.json"

# Instruction to force the agent to call the retrieval tool for every user query
RAG_TOOL_ENFORCE_INSTRUCTION = (
    "IMPORTANT: This is a Retrieval-Augmented Generation (RAG) assistant. For EVERY user "
    "question you MUST call the tool named 'RAG_Document_Search' to fetch relevant passages "
    "from the uploaded document store before producing an answer. Use the retrieved passages "
    "as evidence and prefer them over model-only speculation. If no relevant passages are "
    "found, say so and do not invent facts."
)

@rag_bp.route("/")
def rag():
    """Render the planner template."""
    return render_template("index.html")

# Application logger (console + optional file). Use DEBUG when requested.
logger = logging.getLogger("vibe_app")
logger.setLevel(logging.DEBUG if os.getenv("FLASK_DEBUG") == "1" or os.getenv("DEBUG_LOG") == "1" else logging.INFO)
_fmt = logging.Formatter("%(asctime)s %(levelname)s %(message)s")
_ch = logging.StreamHandler(sys.stdout)
_ch.setFormatter(_fmt)
logger.addHandler(_ch)
try:
    _fh = logging.FileHandler("vibe_app.log")
    _fh.setFormatter(_fmt)
    logger.addHandler(_fh)
except Exception:
    # Best-effort file logging; continue if it fails
    logger.warning("Could not create vibe_app.log (permission or path issue). Continuing without file logging.")

    # ---------------------------------------------------------------------------
    # Module overview
    #
    # This Flask application provides a small web UI and API to interact with
    # different LLM backends (Gemini, OpenAI). It supports:
    # - Setting API keys (persisted to ./api_keys/{llm_name})
    # - Uploading files to be indexed for RAG (stored under ./rag_documents)
    # - Sending prompts to an agent that uses the LLM + a LlamaIndex toolset
    # - Optional "mock" mode to simulate streaming responses for testing
    # - Language filtering/instruction wrapping so the model outputs in a
    #   user-selected target language and supports a 'both' mode (original +
    #   translated output) or 'direct' mode (only the target language)
    #
    # Logging and debugging:
    # - Uses a module logger `logger` which writes to stdout and (optionally)
    #   to `vibe_app.log` in the working directory.
    # - A global error handler returns JSON with a traceback for easier debugging
    #
    # Note: avoid logging secrets (API keys) — we only log high-level request
    # metadata and the selected llm_choice/target_language/response_mode.
    # ---------------------------------------------------------------------------


@rag_bp.before_request
def _log_request_info():
    # Safely log basic request metadata and JSON body (if any). Avoid logging secrets.
    try:
        body = request.get_json(silent=True)
    except Exception:
        body = None
    headers = {
        "User-Agent": request.headers.get("User-Agent"),
        "Content-Type": request.headers.get("Content-Type"),
        "Host": request.headers.get("Host")
    }
    try:
        logger.debug("Incoming request %s %s headers=%s json=%s", request.method, request.path, headers, json.dumps(body) if body is not None else None)
    except Exception:
        logger.debug("Incoming request %s %s (body not JSON-serializable)", request.method, request.path)

@rag_bp.errorhandler(Exception)
def _handle_unhandled_exception(e):
    # Global fall-through for unexpected exceptions: log full traceback and return JSON
    logger.error("Unhandled exception in request", exc_info=True)
    tb = traceback.format_exc()
    return jsonify({"ok": False, "error": "Unhandled server error.", "detail": str(e), "traceback": tb}), 500

def initialize_llm(llm_choice : str) -> None:
    """ Initialize LLM in LlamaIndex friendly way. """
    llm = None

    # The initialize path prepares a LlamaIndex-compatible LLM object. This
    # function currently supports 'gemini' (Google) and can be extended to
    # other providers. It primarily ensures API keys are available and sets
    # Settings.llm so downstream LlamaIndex components can access the model.
    if llm_choice == "gemini":
        api_key = get_environment_api_key(llm_choice)
        if not api_key:
            print("[LLM INIT] Missing Google API key for Gemini.")
            return None

        # Ensure env var for underlying clients
        os.environ["GOOGLE_API_KEY"] = api_key

        try:
            # Use LlamaIndex wrapper
            llm = GoogleGenAI(
                model=GEMINI_MODEL,
                api_key=api_key
            )

        except Exception as e:
            print(f"[LLM INIT] GoogleGenAI constructor failed: {e}")
            return None

        Settings.llm = llm
    
    else:
        # Unknown provider: do nothing. Caller will handle missing LLM.
        pass

def initialize_embedding_model(llm_choice : str):
    """ Initialize embedding model for RAG. """
    embedding_model = None

    if llm_choice == "gemini":
        embedding_model = GoogleGenAIEmbedding(
            model_name = GEMINI_EMBEDDING_MODEL,
            api_key = get_environment_api_key(llm_choice),
            embedding_config = None,
            vertexai_config = None,
            http_options = None,
            debug_config = None,
            embed_batch_size = 1,
            callback_manager = None,
            retries = 3,
            timeout = 10,
            retry_min_seconds = 1,
            retry_max_seconds = 10,
            retry_exponential_base = 2
        )

        Settings.embed_model = embedding_model

    else:
        # If no known embedding model is configured for the requested LLM,
        # leave `Settings.embed_model` unchanged. Downstream code should
        # validate that `Settings.embed_model` is set before attempting
        # to create or load vector indexes.
        pass

def get_vector_index(llm_choice: str):
    """
    Create or load a VectorStoreIndex for the given LLM choice, using a per-LLM
    storage directory and the appropriate embedding model.
    Returns the VectorStoreIndex, or None on failure.
    """
    try:
        index_dir = os.path.join(INDEX_PATH, llm_choice)
        if not Path(index_dir).exists():
            # Fresh index
            vector_store = SimpleVectorStore()
            storage_context = StorageContext.from_defaults(vector_store=vector_store)
            vector_index = VectorStoreIndex(
                nodes=[],
                embed_model=Settings.embed_model,
                use_async=False,
                store_nodes_override=False,
                insert_batch_size=2048,
                storage_context=storage_context,
                callback_manager=None,
                transformations=[]
            )
        else:
            # Load existing index
            storage_context = StorageContext.from_defaults(persist_dir=index_dir)
            vector_index = load_index_from_storage(
                storage_context,
                embed_model=Settings.embed_model,
                use_async=False,
                store_nodes_override=False,
                insert_batch_size=2048,
                callback_manager=None
            )

        return vector_index
    
    except Exception:
        # Return None on any failure so callers can surface an error to the
        # client. Log the exception for debugging.
        logger.exception("Failed to create or load vector index for %s", llm_choice)
        return None

def get_environment_api_key(llm_choice: str) -> str:
    """Get the API key for the selected LLM from ./api_keys/{llm_choice}"""
    key_dir = os.path.join(".", "api_keys")
    key_path = os.path.join(key_dir, llm_choice)
    
    if not os.path.exists(key_path):
        return ""
    try:
        with open(key_path, "r", encoding="utf-8") as f:
            value = f.read().strip()
            return value if value else ""
    
    except OSError:
        return ""

def set_environment_api_key(llm_choice: str, api_key: str):
    """Persist the API key for the selected LLM into ./api_keys/{llm_choice}"""
    key_dir = os.path.join(".", "api_keys")
    os.makedirs(key_dir, exist_ok=True)
    key_path = os.path.join(key_dir, llm_choice)
    # Overwrite any existing value with the new key
    with open(key_path, "w", encoding="utf-8") as f:
        f.write(api_key.strip())

@rag_bp.route("/planner")
def planner():
    """Redirect to the dedicated planner blueprint."""
    return redirect(url_for("planner.planner"))


@rag_bp.route('/notes')
def notes():
    """Render the dedicated notes page."""
    try:
        return render_template('notes.html')
    except Exception:
        return "Notes page template missing.", 200

def _message_to_text(msg) -> str:
    """Extract plain text from message content which can be a str or list of parts."""
    try:
        content = getattr(msg, "content", msg)
        if isinstance(content, str):
            return content
        if isinstance(content, list):
            # Concatenate text-like parts
            parts = []
            for p in content:
                # parts may be dicts or objects with 'type'/'text' fields
                if isinstance(p, dict):
                    t = p.get("text") or p.get("content") or ""
                    if isinstance(t, str):
                        parts.append(t)
                else:
                    t = getattr(p, "text", None)
                    if isinstance(t, str):
                        parts.append(t)
            return "".join(parts)
        return str(content)
    except Exception:
        # If anything goes wrong extracting the message text, return an empty
        # string rather than propagating the exception into the stream.
        logger.exception("Failed to extract text from message object")
        return ""


# ---------------------- Notes helpers and endpoints ------------------------
def _load_notes():
    """Load notes list from NOTES_FILE. Returns list of note dicts."""
    try:
        if not NOTES_FILE.exists():
            return []
        with open(NOTES_FILE, "r", encoding="utf-8") as f:
            data = json.load(f)
            if isinstance(data, list):
                return data
            return []
    except Exception:
        logger.exception("Failed to load notes file")
        return []


def _save_notes(notes):
    """Persist notes list to NOTES_FILE atomically."""
    try:
        NOTES_DIR.mkdir(parents=True, exist_ok=True)
        tmp = NOTES_FILE.with_suffix(".json.tmp")
        with open(tmp, "w", encoding="utf-8") as f:
            json.dump(notes, f, ensure_ascii=False, indent=2)
        os.replace(tmp, NOTES_FILE)
        return True
    except Exception:
        logger.exception("Failed to save notes file")
        return False


@rag_bp.route("/api/notes", methods=["GET"]) 
def list_notes():
    """Return JSON array of saved notes."""
    notes = _load_notes()
    return jsonify({"ok": True, "notes": notes})


@rag_bp.route("/api/notes", methods=["POST"])
def create_or_update_note():
    """Create a new note or update an existing one. Expects JSON {title, content, id?}."""
    data = request.get_json(silent=True) or {}
    title = data.get("title", "")
    content = data.get("content", "")
    note_id = data.get("id")

    notes = _load_notes()
    # If id provided, update
    if note_id is not None:
        for n in notes:
            if n.get("id") == note_id:
                n["title"] = title
                n["content"] = content
                n["updated_at"] = time.strftime("%Y-%m-%dT%H:%M:%SZ", time.gmtime())
                _save_notes(notes)
                return jsonify({"ok": True, "note": n})
        # If not found, fall through to create new

    # Create new note with incremental id
    max_id = max((n.get("id", 0) for n in notes), default=0)
    new = {
        "id": max_id + 1,
        "title": title,
        "content": content,
        "created_at": time.strftime("%Y-%m-%dT%H:%M:%SZ", time.gmtime())
    }
    notes.append(new)
    ok = _save_notes(notes)
    if not ok:
        return jsonify({"ok": False, "error": "Failed to persist note."}), 500
    return jsonify({"ok": True, "note": new})


@rag_bp.route("/api/notes/<int:note_id>", methods=["DELETE"])
def delete_note(note_id: int):
    """Delete the note with the given id."""
    notes = _load_notes()
    new_notes = [n for n in notes if n.get("id") != note_id]
    if len(new_notes) == len(notes):
        return jsonify({"ok": False, "error": "Note not found."}), 404
    ok = _save_notes(new_notes)
    if not ok:
        return jsonify({"ok": False, "error": "Failed to persist notes."}), 500
    return jsonify({"ok": True})

# ---------------------- End notes helpers ---------------------------------

@rag_bp.route("/api/prompt", methods=["POST"])
def receive_prompt():
    """ 
    1. Retrieve frontend variables
    2. Create Langchain agent, return response 
    """
    # Parse JSON body safely. `silent=True` prevents exceptions on invalid JSON.
    data = request.get_json(silent=True) or {}

    # The raw user prompt text (what the user typed in the UI)
    prompt_text = data.get("prompt", "")

    # Which LLM backend the user selected (e.g., 'gemini')
    llm_choice = data.get("llm_choice", "gemini")

    # Accept either 'target_language' or 'language' from the frontend for
    # backward compatibility with different client payloads.
    target_language = data.get("target_language", "") or data.get("language", "")

    # response_mode controls how we instruct the LLM:
    # - 'direct': reply only in the target language
    # - 'both': include original answer then a translated version separated
    #           by a marker (e.g. '---TRANSLATION (Spanish)---')
    response_mode = data.get("response_mode", "direct")

    # Basic validation: ensure caller selected an LLM and we have an API key

    if llm_choice == "":
        return jsonify({"ok": False, "error": "SNO: no LLM selected."}), 400

    api_key: str = get_environment_api_key(llm_choice)
    if api_key == "":
        return jsonify({"ok": False, "error": "NO API key set."}), 400

    # Wrap the handler in try/except to capture unexpected errors and print
    # a stacktrace to the server logs for diagnosis rather than returning a
    # generic 500 without information.
    try:
        # Log key request parameters (do NOT log raw API keys)
        logger.debug("/api/prompt called; llm_choice=%s target_language=%s response_mode=%s", llm_choice, target_language, response_mode)
        if llm_choice == "gemini":
            langchain_llm: BaseChatModel = ChatGoogleGenerativeAI(
                model=GEMINI_MODEL,
                google_api_key=api_key,
                temperature = 0.1
            )
        else:
            return jsonify({"ok": False, "error": "Invalid LLM selected."}), 400

        global is_first_llm_run
        if is_first_llm_run == 1:
            print("Initializing LLM")
            initialize_llm(llm_choice)
            is_first_llm_run = 0

        global is_first_embed_run
        if is_first_embed_run == 1:
            print("Initializing embedding model")
            initialize_embedding_model(llm_choice)
            is_first_embed_run = 0

        # Obtain VectorStoreIndex for this llm_choice (for file ingestion/RAG)
        vector_index = get_vector_index(llm_choice)
        if vector_index is None:
            return jsonify({"ok": False, "error": "Unable to create or load index."}), 400

        print("Vector index made")

        # Create a query engine and expose it as a text-returning tool
        query_engine = vector_index.as_query_engine(
            similarity_top_k=10,    # how many docs are returned to the synthesizer
            fetch_k=50,             # how many are fetched from the store before reranking
            similarity_cutoff=0.1,  # include lower-similarity candidates (raise if too noisy)
            response_mode="compact" # try 'compact' or 'tree_summarize' depending on your llama-index version
        )
        
        # Monkey-patch the created query_engine to add lightweight debug logging
        # for retrievals without changing its type (avoids pydantic type checks).
        def _add_debug_wrappers(qe):
            def make_wrapper(orig_fn):
                def wrapper(*args, **kwargs):
                    try:
                        print("[RAG DEBUG] Running query with args=", args, "kwargs=", {k:v for k,v in kwargs.items() if k!='messages'})
                    except Exception:
                        pass
                    res = orig_fn(*args, **kwargs)
                    try:
                        nodes = None
                        if hasattr(res, 'source_nodes'):
                            nodes = res.source_nodes
                        elif isinstance(res, dict) and 'source_nodes' in res:
                            nodes = res['source_nodes']
                        elif hasattr(res, 'extra_info') and isinstance(res.extra_info, dict):
                            nodes = res.extra_info.get('source_nodes')

                        if nodes:
                            print(f"[RAG DEBUG] Retrieved {len(nodes)} nodes:")
                            for i, n in enumerate(nodes, start=1):
                                try:
                                    meta = getattr(n, 'metadata', None) or (n.get('metadata') if isinstance(n, dict) else None)
                                    text = getattr(n, 'text', None) or (n.get('text') if isinstance(n, dict) else None)
                                    score = getattr(n, 'score', None) if hasattr(n, 'score') else (n.get('score') if isinstance(n, dict) else None)
                                    text_snippet = ((text or '')[:200]).replace('\n', ' ')
                                    print("  %s. score=%s metadata=%s text_snippet=%s" % (i, score, meta, text_snippet))
                                except Exception as e:
                                    print(f"  {i}. <failed to print node>: {e}")
                        else:
                            print("[RAG DEBUG] No source_nodes found on result; result type:", type(res))
                    except Exception as e:
                        print("[RAG DEBUG] Error while extracting nodes:", e)
                    return res
                return wrapper

            # Try to wrap common methods without changing the object's type
            for name in ("query", "run", "__call__"):
                orig = getattr(qe, name, None)
                if callable(orig):
                    try:
                        setattr(qe, name, make_wrapper(orig))
                    except Exception as e:
                        print(f"[RAG DEBUG] Failed to patch method {name}: {e}")

        try:
            _add_debug_wrappers(query_engine)
            print("Query engine made and monkey-patched with debug wrappers")
        except Exception as e:
            print("[RAG DEBUG] Failed to add debug wrappers:", e)

        # Optional retrieval-only debug (no agent) to inspect what the retriever returns.
        # Enable by setting environment variable RAG_RETRIEVE_DEBUG=1. This may
        # invoke the retriever and/or the LLM depending on the engine implementation.
        if os.getenv("RAG_RETRIEVE_DEBUG", "0") == "1":
            test_query = prompt_text or "test retrieval"
            tried = []
            methods = ["retrieve", "retrieve_nodes", "get_relevant_documents", "query", "run", "__call__"]
            for m in methods:
                fn = getattr(query_engine, m, None)
                if callable(fn):
                    try:
                        print(f"[RAG RETRIEVE DEBUG] Calling method: {m}()")
                        # Some methods expect different signatures; try common ones.
                        try:
                            r = fn(test_query)
                        except TypeError:
                            # Try passing as kwargs
                            r = fn(query=test_query)
                        print(f"[RAG RETRIEVE DEBUG] Result type from {m}:", type(r))
                        # Try to extract source_nodes or documents
                        nodes = None
                        if hasattr(r, 'source_nodes'):
                            nodes = r.source_nodes
                        elif isinstance(r, dict) and 'source_nodes' in r:
                            nodes = r['source_nodes']
                        elif hasattr(r, 'documents'):
                            nodes = r.documents
                        elif isinstance(r, list) and r and hasattr(r[0], 'metadata'):
                            nodes = r

                        if nodes:
                            print(f"[RAG RETRIEVE DEBUG] {m} returned {len(nodes)} nodes:")
                            for i, n in enumerate(nodes, start=1):
                                try:
                                    meta = getattr(n, 'metadata', None) or (n.get('metadata') if isinstance(n, dict) else None)
                                    text = getattr(n, 'text', None) or (n.get('text') if isinstance(n, dict) else None)
                                    score = getattr(n, 'score', None) if hasattr(n, 'score') else (n.get('score') if isinstance(n, dict) else None)
                                    text_snippet = ((text or '')[:300]).replace('\n', ' ')
                                    print("  %s. score=%s metadata=%s text_snippet=%s" % (i, score, meta, text_snippet))
                                except Exception as e:
                                    print(f"  {i}. <failed to print node>: {e}")
                        else:
                            print(f"[RAG RETRIEVE DEBUG] {m} did not return nodes (raw result: {r})")
                        tried.append(m)
                    except Exception as e:
                        print(f"[RAG RETRIEVE DEBUG] call {m} failed: {e}")
            if not tried:
                print("[RAG RETRIEVE DEBUG] No retrieval methods found on query_engine to call.")

        query_tool_config = IndexToolConfig(
            query_engine = query_engine,
            name = "RAG_Document_Search",
            description=(
                "Use this tool to search the uploaded document store for passages relevant to the user's question. "
                "ALWAYS call this tool for factual, document-based, or specific queries about uploaded content. "
                "Return retrieved passages and metadata and use them as evidence when composing answers."
            ),
        )

        query_engine_tool = LlamaIndexTool.from_tool_config(query_tool_config)
        print("Tool made")

        # Wrap the tool object's callable entrypoints so we log when the agent
        # invokes the tool. This captures agent -> tool calls that may bypass
        # direct query_engine method invocations.
        try:
            def _wrap_tool_invocation(tool_obj):
                # Try common method names and wrap the first callable we find.
                for method_name in ("run", "__call__", "invoke", "execute", "query"):
                    orig = getattr(tool_obj, method_name, None)
                    if callable(orig):
                        def make_wrapper(orig_fn, method_name=method_name):
                            def wrapper(*args, **kwargs):
                                try:
                                    print(f"[RAG TOOL DEBUG] Tool invoked via {method_name} args=", args, "kwargs=", {k:v for k,v in kwargs.items() if k!='messages'})
                                except Exception:
                                    pass
                                res = orig_fn(*args, **kwargs)
                                # Try to extract retrieved nodes/documents from result
                                try:
                                    nodes = None
                                    if hasattr(res, 'source_nodes'):
                                        nodes = res.source_nodes
                                    elif isinstance(res, dict) and 'source_nodes' in res:
                                        nodes = res['source_nodes']
                                    elif hasattr(res, 'documents'):
                                        nodes = res.documents
                                    elif isinstance(res, list) and res and hasattr(res[0], 'metadata'):
                                        nodes = res

                                    if nodes:
                                        print(f"[RAG TOOL DEBUG] {method_name} returned {len(nodes)} nodes:")
                                        for i, n in enumerate(nodes, start=1):
                                            try:
                                                meta = getattr(n, 'metadata', None) or (n.get('metadata') if isinstance(n, dict) else None)
                                                text = getattr(n, 'text', None) or (n.get('text') if isinstance(n, dict) else None)
                                                score = getattr(n, 'score', None) if hasattr(n, 'score') else (n.get('score') if isinstance(n, dict) else None)
                                                text_snippet = ((text or '')[:300]).replace('\n', ' ')
                                                print("  %s. score=%s metadata=%s text_snippet=%s" % (i, score, meta, text_snippet))
                                            except Exception as e:
                                                print(f"  {i}. <failed to print node>: {e}")
                                    else:
                                        print(f"[RAG TOOL DEBUG] {method_name} did not return nodes; raw result type: {type(res)}")
                                except Exception as e:
                                    print("[RAG TOOL DEBUG] Error while extracting nodes from tool result:", e)
                                return res
                            return wrapper
                        try:
                            setattr(tool_obj, method_name, make_wrapper(orig))
                            print(f"[RAG TOOL DEBUG] Wrapped tool method: {method_name}")
                        except Exception as e:
                            print(f"[RAG TOOL DEBUG] Failed to wrap tool method {method_name}: {e}")
                        break

            _wrap_tool_invocation(query_engine_tool)
        
        except Exception as e:
            print("[RAG TOOL DEBUG] Exception while attempting to wrap tool invocations:", e)

        tools = [query_engine_tool]

        langchain_agent: CompiledStateGraph = create_react_agent(
            model=langchain_llm,
            tools=tools
        )
        print("Agent made")
        logger.debug("Agent created for llm_choice=%s", llm_choice)

        # Build a system-level instruction (enforcement + optional language rules)
        # so the model treats these as instructions rather than content to echo/translate.
        try:
            system_parts = [RAG_TOOL_ENFORCE_INSTRUCTION]
            lang_name = None

            # If a target language is requested, append instructions to the system message
            if target_language:
                language_map = {
                    "en": "English",
                    "es": "Spanish",
                    "fr": "French",
                    "de": "German",
                    "zh": "Chinese (Mandarin)",
                    "hi": "Hindi",
                    "ar": "Arabic",
                    "pt": "Portuguese",
                    "ru": "Russian",
                    "it": "Italian",
                    "ja": "Japanese",
                    "ko": "Korean",
                    "tr": "Turkish",
                    "nl": "Dutch",
                    "sv": "Swedish",
                    "pl": "Polish",
                    "vi": "Vietnamese",
                    "th": "Thai",
                    "id": "Indonesian",
                    "bn": "Bengali",
                    "ur": "Urdu",
                    "fa": "Persian",
                    "he": "Hebrew",
                    "ro": "Romanian",
                    "cs": "Czech",
                    "el": "Greek",
                    "hu": "Hungarian",
                    "no": "Norwegian",
                    "sk": "Slovak"
                }
                lang_name = language_map.get(target_language, target_language)

                if response_mode == "both":
                    system_parts.append(
                        f"Provide a complete answer to the user's question. After the full answer, "
                        f"insert a line that says '---TRANSLATION ({lang_name})---' and then provide a "
                        f"translation of the full answer into {lang_name}. Do not include any additional commentary."
                    )
                else:
                    system_parts.append(f"Please respond ONLY in {lang_name}. All output should be in {lang_name}.")

            # Build the final structured messages list for the agent
            system_content = "\n\n".join(system_parts)
            messages = [
                {"role": "system", "content": system_content},
                {"role": "user", "content": prompt_text}
            ]
        except Exception:
            # If anything goes wrong, fall back to a simple user message
            messages = [{"role": "user", "content": prompt_text}]
        
        # Mock LLM mode: allow testing without external API calls. Enable by setting
        # the environment variable MOCK_LLM=1 or including {"mock": true} in the POST body.
        mock_mode = (os.getenv("MOCK_LLM", "0") == "1") or bool(data.get("mock", False))
        
        if mock_mode:
            logger.info("Mock LLM mode active - streaming a simulated response")
            # Use a small generator to simulate streaming chunks
            def mock_gen():
                yield ""
                time.sleep(0.01)
                # main answer (shortened)
                main = f"[MOCK ANSWER] Responding to: {prompt_text}"
                for i in range(0, len(main), 100):
                    yield main[i:i+100]
                    time.sleep(0.01)
                if response_mode == "both":
                    marker = f"\n---TRANSLATION ({lang_name if 'lang_name' in locals() else target_language})---\n"
                    for i in range(0, len(marker), 100):
                        yield marker[i:i+100]
                        time.sleep(0.01)
                    trans = f"[MOCK TRANSLATION into {lang_name if 'lang_name' in locals() else target_language}]"
                    for i in range(0, len(trans), 100):
                        yield trans[i:i+100]
                        time.sleep(0.01)
                yield "\n"

            return Response(stream_with_context(mock_gen()), mimetype="text/plain")
        
        @stream_with_context
        def generate():
            # Optional: small preamble so client can clear UI
            yield ""

            for step in langchain_agent.stream({"messages": messages}, stream_mode="values"):
                # Optional stream-step debugging: when RAG_DEBUG=1, print step summary
                if os.getenv("RAG_DEBUG", "0") == "1":
                    try:
                        # Avoid printing large structures; show keys and small repr
                        keys = list(step.keys()) if isinstance(step, dict) else []
                        print(f"[RAG STREAM DEBUG] step keys={keys} type={type(step)}")
                    except Exception:
                        pass

                msg = step["messages"][-1]

                # Only yield if the message is from the AI (not Human)
                if isinstance(msg, AIMessage):
                    text = _message_to_text(msg)
                    if text:
                        yield text

            # Optionally end with a newline
            yield "\n"

        # Normal return: stream generator
        return Response(generate(), mimetype="text/plain")
    
    except Exception as e:
        # Print full traceback to server logs for debugging
        traceback.print_exc()
        # Return JSON error so frontend can display a useful message
        return jsonify({"ok": False, "error": "Internal server error.", "detail": str(e)}), 500

@rag_bp.route("/api/set-api-key", methods=["POST"])
def set_api_key():
    """ Set the API key for the selected LLM """
    data = request.get_json(silent=True) or {}
    llm_choice = data.get("llm_choice", "gemini")
    api_key = data.get("api_key", "")
  
    if llm_choice == "":
        return jsonify({"ok": False, "error": "SNO: no LLM choice provided."}), 400
        
    if api_key == "":
        return jsonify({"ok": False, "error": "SNO: no API key provided."}), 400
    
    # Set the environment variable based on LLM choice
    set_environment_api_key(llm_choice, api_key)
    
    return jsonify({"ok": True, "message": f"API key set for {llm_choice}"})


@rag_bp.route("/api/upload-files", methods=["POST"])
def upload_files():
    """ Handle file uploads and print file information """
    data = request.get_json(silent=True) or {}
    files = data.get("file_paths", [])
    llm_choice = data.get("llmChoice", "gemini")
    
    api_key: str = get_environment_api_key(llm_choice)
    if api_key == "":
        return jsonify({"ok": False, "error": "NO API key set."}), 400

    global is_first_embed_run
    if is_first_embed_run == 1:
        print("Initializing embedding model")
        initialize_embedding_model(llm_choice)
        is_first_embed_run = 0

    # Obtain VectorStoreIndex for this llm_choice (for file ingestion/RAG)
    vector_index = get_vector_index(llm_choice)
    if vector_index is None:
        return jsonify({"ok": False, "error": "Unable to create or load index."}), 400

    # Add to RAG index as PDFs are parsed
    any_inserted = False
    splitter = SentenceSplitter(chunk_size=1200, chunk_overlap=200)

    for file_info in files:
        content_b64 = file_info.get('content', '')
        file_name = file_info.get('name', '')
        file_type = file_info.get('type', '')
        extracted_text = ""

        if not content_b64:
            print(f"[FILE_UPLOAD] No content found for {file_name}")
            continue
        
        if file_type == "application/pdf":
            try:
                with pdfplumber.open(BytesIO(base64.b64decode(content_b64))) as pdf:
                    for page_number, page in enumerate(pdf.pages, 1):
                        text = page.extract_text() or ""
                        extracted_text += f"\n\n--- Page {page_number} ---\n{text}\n"
                        
                        '''
                        # Tried to do images, didn't go well
                        # Discussed more in the paper
                        # Extract the images from the PDF
                        for image_idx, image in enumerate(page.images):
                            try:
                                print(f"[FILE_UPLOAD] Image: {image}")
                                image_bounding_box = (image["x0"], page.height - image["y1"], image["x1"], page.height - image["y0"])
                                cropped_image = page.crop(image_bounding_box).to_image(resolution=300)
                               
                                # Convert cropped_image to bytes and then to base64
                                cropped_bytes = BytesIO()
                                cropped_image.save(cropped_bytes, format="PNG")
                                cropped_bytes.seek(0)
                                cropped_base64 = base64.b64encode(cropped_bytes.read()).decode("utf-8")
                               
                                # Create ImageDocument for multimodal RAG
                                image_document = ImageDocument(
                                    image=cropped_base64,
                                    image_mimetype="image/png",
                                    metadata={
                                        "file_name": file_name,
                                        "page_num": page_number,
                                        "image_idx": image_idx,
                                        "source_pdf": file_name,
                                        "file_type": file_type
                                    }
                                )
                               
                                # Insert image document into vector index
                                print(f"[RAG] Inserting image document from {file_name}, page {page_number}, image {image_idx}...")
                                vector_index.insert_nodes([image_document])
                                any_inserted = True


                            except Exception as img_error:
                                print(f"[IMAGE UPLOAD] Error with image detection: {img_error}. Continuing.")
                                continue

                        # Extract the images from the PDF
                        for image_index, image in enumerate(page.images):
                            try:
                                print(f"[FILE_UPLOAD] Image: {image}")
                                image_bounding_box = (image["x0"], page.height - image["y1"], image["x1"], page.height - image["y0"])
                                cropped_image = page.crop(image_bounding_box).to_image(resolution=300)
                                print(type(cropped_image))
                                image_document = ImageDocument(
                                    image = Image.fromarray(np.array(cropped_image.original)),
                                    metadata={
                                        "file_name": file_name,
                                        "page_num": page_number,
                                        "image_idx": image_index,
                                        "source_pdf": file_name
                                    }
                                )
                                # TODO: add image to RAG database

                            except Exception as img_error:
                                print(f"[IMAGE UPLOAD] Error with image detection: {img_error}. Continuing.")
                                continue
                        '''

                        # Convert tables to text descriptions
                        for table in page.extract_tables():
                            df = pd.DataFrame(table[1:], columns=table[0])
                            text += f"\n\nTable:\n{df.to_string()}\n"
                            extracted_text += text
                            text = ""

                if extracted_text.strip():
                    doc = Document(text=extracted_text, metadata={"file_name": file_name, "file_type": file_type})
                    nodes = splitter.get_nodes_from_documents([doc])
                    print(f"[RAG] Inserting {len(nodes)} nodes for {file_name}...")
                    vector_index.insert_nodes(nodes)
                    print(f"[RAG] Inserted nodes for {file_name}")
                    any_inserted = True
        
            except Exception as e:
                print(f"[FILE_UPLOAD] Error extracting text from {file_name}: {e}")
        
        elif file_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
            # DOCX parsing (text + tables)
            try:
                docx_obj = DocxDocument(BytesIO(base64.b64decode(content_b64)))
                for para in docx_obj.paragraphs:
                    extracted_text += para.text + "\n"

                # Tables
                for table_idx, table in enumerate(docx_obj.tables, start=1):
                    try:
                        rows = [[cell.text for cell in row.cells] for row in table.rows]
                        df = pd.DataFrame(rows)
                        extracted_text += f"\n\n--- Table {table_idx} ---\n{df.to_string(index=False, header=False)}\n"
                    except Exception as tbl_err:
                        print(f"[DOCX TABLE] Error parsing table in {file_name}: {tbl_err}")
                
                if extracted_text.strip():
                    doc = Document(text=extracted_text, metadata={"file_name": file_name, "file_type": file_type})
                    nodes = splitter.get_nodes_from_documents([doc])
                    print(f"[RAG] Inserting {len(nodes)} nodes for {file_name} (docx)...")
                    vector_index.insert_nodes(nodes)
                    print(f"[RAG] Inserted nodes for {file_name} (docx)")
                    any_inserted = True
            except Exception as e:
                print(f"[FILE_UPLOAD] Error extracting text from {file_name} (docx): {e}")

        elif file_type in ("application/vnd.openxmlformats-officedocument.presentationml.presentation", "application/vnd.ms-powerpoint"):
            # PPTX parsing (text from shapes + tables)
            try:
                prs = Presentation(BytesIO(base64.b64decode(content_b64)))
                for slide_num, slide in enumerate(prs.slides, start=1):
                    slide_text_parts = []
                    for shape in slide.shapes:
                        try:
                            if hasattr(shape, "text") and shape.text:
                                slide_text_parts.append(shape.text)
                        except Exception:
                            # Some shapes may throw on access; skip them
                            continue

                        # Table extraction (if present)
                        try:
                            if getattr(shape, "has_table", False):
                                table = shape.table
                                rows = []
                                for r in table.rows:
                                    cells = [c.text for c in r.cells]
                                    rows.append(cells)
                                df = pd.DataFrame(rows)
                                slide_text_parts.append(f"\nTable:\n{df.to_string(index=False, header=False)}\n")
                        except Exception:
                            continue

                    if slide_text_parts:
                        extracted_text += f"\n\n--- Slide {slide_num} ---\n" + "\n".join(slide_text_parts) + "\n"

                if extracted_text.strip():
                    doc = Document(text=extracted_text, metadata={"file_name": file_name, "file_type": file_type})
                    nodes = splitter.get_nodes_from_documents([doc])
                    print(f"[RAG] Inserting {len(nodes)} nodes for {file_name} (pptx)...")
                    vector_index.insert_nodes(nodes)
                    print(f"[RAG] Inserted nodes for {file_name} (pptx)")
                    any_inserted = True
            except Exception as e:
                print(f"[FILE_UPLOAD] Error extracting text from {file_name} (pptx): {e}")

        elif file_type in ("text/plain", "text/markdown"):
            # Plain text / Markdown
            try:
                text = base64.b64decode(content_b64).decode("utf-8", errors="replace")
                extracted_text += text
                if extracted_text.strip():
                    doc = Document(text=extracted_text, metadata={"file_name": file_name, "file_type": file_type})
                    nodes = splitter.get_nodes_from_documents([doc])
                    print(f"[RAG] Inserting {len(nodes)} nodes for {file_name} (text)...")
                    vector_index.insert_nodes(nodes)
                    print(f"[RAG] Inserted nodes for {file_name} (text)")
                    any_inserted = True
            except Exception as e:
                print(f"[FILE_UPLOAD] Error extracting text from {file_name} (text): {e}")

        elif file_type == "text/html":
            # HTML parsing: extract visible text and simple table conversion
            try:
                html = base64.b64decode(content_b64).decode("utf-8", errors="replace")
                soup = BeautifulSoup(html, "lxml")
                # Remove script/style
                for s in soup(["script", "style"]):
                    s.extract()

                # Extract tables
                for t_idx, table in enumerate(soup.find_all("table"), start=1):
                    try:
                        rows = []
                        for tr in table.find_all("tr"):
                            cells = [td.get_text(separator=" ", strip=True) for td in tr.find_all(["td", "th"])]
                            rows.append(cells)
                        if rows:
                            df = pd.DataFrame(rows)
                            extracted_text += f"\n\n--- HTML Table {t_idx} ---\n{df.to_string(index=False, header=False)}\n"
                    except Exception as et:
                        print(f"[HTML TABLE] Error parsing table in {file_name}: {et}")

                # Visible text
                visible = soup.get_text(separator="\n", strip=True)
                extracted_text += "\n" + visible

                if extracted_text.strip():
                    doc = Document(text=extracted_text, metadata={"file_name": file_name, "file_type": file_type})
                    nodes = splitter.get_nodes_from_documents([doc])
                    print(f"[RAG] Inserting {len(nodes)} nodes for {file_name} (html)...")
                    vector_index.insert_nodes(nodes)
                    print(f"[RAG] Inserted nodes for {file_name} (html)")
                    any_inserted = True

            except Exception as e:
                print(f"[FILE_UPLOAD] Error extracting text from {file_name} (html): {e}")

        else:
            print(f"[FILE_UPLOAD] Unsupported file type or empty content: name={file_name}, type={file_type}")

    # Persist once after all inserts (persist to per-LLM directory)
    if any_inserted:
        print("Persisting index")
        index_dir = os.path.join(INDEX_PATH, llm_choice)
        
        # Make sure index directory exists
        if not os.path.exists(index_dir):
            print("Here")
            os.makedirs(index_dir, exist_ok=True)
        
        vector_index.storage_context.persist(Path(index_dir))
    
    return jsonify({"ok": True, "message": f"Received {len(files)} file(s)"})


@rag_bp.route("/api/scrape", methods=["POST"])
def scrape_url():
    """Fetch a URL, extract visible text, and optionally insert into the RAG index.

    POST JSON parameters:
    - url: string (required) - http(s) URL to fetch
    - insert: bool (optional) - if true, insert extracted content into vector index
    - llm_choice: string (required if insert=true) - which LLM/index to insert into
    """
    data = request.get_json(silent=True) or {}
    url = (data.get("url") or "").strip()
    insert = bool(data.get("insert", False))
    llm_choice = data.get("llm_choice", "gemini")

    if not url:
        return jsonify({"ok": False, "error": "No URL provided."}), 400

    if not (url.startswith("http://") or url.startswith("https://")):
        return jsonify({"ok": False, "error": "Only http:// and https:// URLs are supported."}), 400

    try:
        req = urllib.request.Request(url, headers={
            "User-Agent": "vibe-scraper/1.0 (+https://example.local)"
        })
        with urllib.request.urlopen(req, timeout=15) as resp:
            raw = resp.read()
            # Try to decode using detected encoding or utf-8 fallback
            try:
                charset = resp.headers.get_content_charset() or "utf-8"
            except Exception:
                charset = "utf-8"
            try:
                html = raw.decode(charset, errors="replace")
            except Exception:
                html = raw.decode("utf-8", errors="replace")

    except urllib.error.HTTPError as e:
        logger.exception("HTTP error fetching %s", url)
        return jsonify({"ok": False, "error": f"HTTP error: {e.code} fetching URL."}), 502
    except urllib.error.URLError as e:
        logger.exception("URL error fetching %s", url)
        return jsonify({"ok": False, "error": f"URL fetch failed: {e.reason}"}), 502
    except Exception as e:
        logger.exception("Unexpected error fetching %s", url)
        return jsonify({"ok": False, "error": "Failed to fetch URL."}), 502

    try:
        soup = BeautifulSoup(html, "lxml")
        for s in soup(["script", "style", "noscript"]):
            s.extract()
        text = soup.get_text(separator="\n", strip=True)
    except Exception:
        logger.exception("Failed to parse HTML for %s", url)
        return jsonify({"ok": False, "error": "Failed to parse HTML."}), 500

    inserted = False
    if insert:
        if not llm_choice:
            return jsonify({"ok": False, "error": "llm_choice required when insert=true"}), 400

        # Ensure embedding model initialized for this LLM
        global is_first_embed_run
        if is_first_embed_run == 1:
            initialize_embedding_model(llm_choice)
            is_first_embed_run = 0

        vector_index = get_vector_index(llm_choice)
        if vector_index is None:
            return jsonify({"ok": False, "error": "Unable to create or load index for insertion."}), 500

        try:
            doc = Document(text=text, metadata={"source_url": url})
            splitter = SentenceSplitter(chunk_size=1200, chunk_overlap=200)
            nodes = splitter.get_nodes_from_documents([doc])
            logger.info("Inserting %d nodes from scraped URL %s into index %s", len(nodes), url, llm_choice)
            vector_index.insert_nodes(nodes)
            # Persist per-LLM index directory
            index_dir = os.path.join(INDEX_PATH, llm_choice)
            os.makedirs(index_dir, exist_ok=True)
            vector_index.storage_context.persist(Path(index_dir))
            inserted = True
        except Exception:
            logger.exception("Failed to insert scraped content into index for %s", url)
            return jsonify({"ok": False, "error": "Failed to insert into index."}), 500

    # Return the extracted text (possibly truncated) and insertion status
    max_return = 20000
    short = text if len(text) <= max_return else text[:max_return] + "\n\n...[truncated]..."
    return jsonify({"ok": True, "url": url, "text": short, "inserted": inserted})
